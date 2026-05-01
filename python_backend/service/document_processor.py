"""
Document Processing Service with RAG using Pinecone
Handles document upload, chunking, embedding, and vector storage
Supports PDF, DOCX, PPT, and images with S3 storage integration
"""
import os
import uuid
import asyncio
import re
import hashlib
from typing import List, Dict, Optional, Any, Generator
from datetime import datetime, timedelta
import json
import logging

from pinecone import PineconeAsyncio, ServerlessSpec
from openai import AsyncOpenAI
from pydantic import BaseModel

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from io import BytesIO

from configs import config
from service.database import get_database_service
from service.document_heuristics import dedup_text_snippets, filter_by_relevance
from service.s3_service import get_s3_service
from agents.file_parser_agent import get_file_parser_agent

logger = logging.getLogger(__name__)

# File size limits (in bytes)
MAX_FILE_SIZE = int(os.getenv("MAX_FILE_SIZE", str(getattr(config, 'MAX_FILE_SIZE', 10485760))))  # Default 10MB
MAX_IMAGE_SIZE = int(os.getenv("MAX_IMAGE_SIZE", str(getattr(config, 'MAX_IMAGE_SIZE', 5242880))))  # Default 5MB


class DocumentChunk(BaseModel, extra="allow"):
    """Represents a chunk of a document (small rows; agents pull adjacent rows for wider context)."""
    chunk_index: int
    content: str
    metadata: Dict[str, Any]
    document_id: Optional[str] = None
    chunk_id: Optional[str] = None
    content_hash: Optional[str] = None


class DocumentProcessor:
    """
    Handles document processing, chunking, embedding, and RAG operations
    Memory-efficient for large files with streaming chunking
    """
    
    def __init__(self, model_name: Optional[str] = None, api_key: Optional[str] = None):
        """Initialize document processor"""
        # Initialize Pinecone
        self.pinecone_api_key = os.getenv("PINECONE_API_KEY")
        # if not self.pinecone_api_key:
        #     raise ValueError("PINECONE_API_KEY environment variable is required")
        
        self.pinecone = PineconeAsyncio(api_key=self.pinecone_api_key)
        
        # Initialize OpenAI for embeddings
        self.openai_api_key = os.getenv("OPENAI_API_KEY") or config.OPENAI_API_KEY
        self.openai_client = AsyncOpenAI(api_key=self.openai_api_key) if self.openai_api_key else None
        
        # Initialize services
        self.db_service = get_database_service()
        self.s3_service = get_s3_service()
        self.file_parser_agent = get_file_parser_agent()
        
        # Default expiration time (2 hours for temporary storage)
        self.default_expiration_hours = int(os.getenv("PINECONE_INDEX_EXPIRATION_HOURS", "2"))
        self.embedding_dimension = 1536  # OpenAI text-embedding-3-small dimension

    @staticmethod
    def _content_hash(text: str) -> str:
        return hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()

    @staticmethod
    def _stable_chunk_id(document_id: str, chunk_index: int, content_hash: str) -> str:
        return f"{document_id}:{chunk_index}:{content_hash[:16]}"

    async def _generate_embedding(self, text: str) -> List[float]:
        """Embed text via OpenAI (matches Pinecone index dimension=1536)."""
        if not self.openai_client:
            raise RuntimeError("OPENAI_API_KEY required for query embeddings")
        resp = await self.openai_client.embeddings.create(
            model="text-embedding-3-small",
            input=text[:8000],  # OpenAI input cap is generous; keep bounded
        )
        return resp.data[0].embedding
    
    def validate_file_size(self, file_content: bytes, file_type: str) -> None:
        """Validate file size against limits"""
        file_size = len(file_content)
        max_size = MAX_IMAGE_SIZE if file_type.startswith("image/") else MAX_FILE_SIZE
        
        if file_size > max_size:
            raise ValueError(
                f"File size ({file_size / 1024 / 1024:.2f}MB) exceeds maximum allowed size "
                f"({max_size / 1024 / 1024:.2f}MB)"
            )
    
    async def process_document(
        self,
        file_content: bytes,
        file_name: str,
        file_type: str,
        user_id: str,
        chunk_size: Optional[int] = None,
        overlap: Optional[int] = None,
        document_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Process uploaded document: extract text, chunk (small word windows), embed, Pinecone, DB `document_chunks`.
        Also saves to S3 if enabled for re-ingestion.
        """
        document_id = document_id or str(uuid.uuid4())
        if chunk_size is None:
            chunk_size = int(getattr(config, "DOCUMENT_CHUNK_WORDS", 1000) or 1000)
        if overlap is None:
            overlap = int(getattr(config, "DOCUMENT_CHUNK_OVERLAP", 200) or 200)
        
        # Validate file size
        self.validate_file_size(file_content, file_type)
        
        # Save to S3 if enabled (for re-ingestion after expiration)
        s3_key = None
        if self.s3_service.enabled:
            s3_key = await self.s3_service.upload_file(
                file_content=file_content,
                user_id=user_id,
                document_id=document_id,
                file_name=file_name,
                content_type=file_type,
                metadata={"uploaded_at": datetime.now().isoformat()}
            )
        
        # Extract text from document
        text_content = await self._extract_text(file_content, file_type, file_name)
        
        if not text_content or not text_content.strip():
            raise ValueError(f"Could not extract text from {file_type} file")
        
        # Check if text is nonsensical (for PDF/DOCX)
        if file_type in ["application/pdf", "application/msword", 
                         "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            if self._is_text_nonsensical(text_content):
                logger.info(f"Text appears nonsensical, using file_parser_agent for {file_name}")
                # Use file parser agent for better extraction
                parsed_result = await self.file_parser_agent.parse_binary_content(
                    content_data=file_content,
                    media_type=file_type,
                    filename=file_name,
                )
                if parsed_result and hasattr(parsed_result, 'extracted_text'):
                    text_content = parsed_result.extracted_text
                elif isinstance(parsed_result, dict) and parsed_result.get('extracted_text'):
                    text_content = parsed_result['extracted_text']

        max_words = int(getattr(config, "MAX_EXTRACTED_WORDS", 200000) or 200000)
        tw = text_content.split()
        if len(tw) > max_words:
            text_content = " ".join(tw[:max_words])
        full_content_hash = self._content_hash(text_content)
        raw_chunks = list(
            self._chunk_document_streaming(text_content, chunk_size=chunk_size, overlap=overlap)
        )
        chunks: List[DocumentChunk] = []
        for ch in raw_chunks:
            h = self._content_hash(ch.content)
            cid = self._stable_chunk_id(document_id, ch.chunk_index, h)
            md = {
                **(ch.metadata or {}),
                "content_hash": h,
                "chunk_id": cid,
                "document_id": document_id,
            }
            chunks.append(
                DocumentChunk(
                    chunk_index=ch.chunk_index,
                    content=ch.content,
                    metadata=md,
                    document_id=document_id,
                    chunk_id=cid,
                    content_hash=h,
                )
            )

        index_name = f"doc-{document_id}-{int(datetime.now().timestamp())}"
        await self._create_pinecone_index(index_name)
        index = await self.pinecone.get_index(index_name)

        # Embed all chunks (batched) then upsert as vector records to match the
        # 1536-d index. Upsert in slices to stay under Pinecone's per-call cap.
        if not self.openai_client:
            raise RuntimeError("OPENAI_API_KEY required for chunk embeddings")
        emb_inputs = [c.content[:8000] for c in chunks]
        embeddings: List[List[float]] = []
        for i in range(0, len(emb_inputs), 100):
            batch = emb_inputs[i : i + 100]
            resp = await self.openai_client.embeddings.create(
                model="text-embedding-3-small", input=batch
            )
            embeddings.extend(d.embedding for d in resp.data)

        vectors: List[Dict[str, Any]] = []
        for ch, vec in zip(chunks, embeddings):
            rid = f"{document_id}_{ch.chunk_index}".replace(" ", "_")[:512]
            vectors.append(
                {
                    "id": rid,
                    "values": vec,
                    "metadata": {
                        "document_id": document_id,
                        "chunk_id": ch.chunk_id,
                        "content_hash": ch.content_hash,
                        "chunk_index": ch.chunk_index,
                        "file_name": file_name,
                        "start_word": ch.metadata.get("start_word", 0),
                        "end_word": ch.metadata.get("end_word", 0),
                        # Query responses need snippet text; bounded for metadata limits
                        "content": ch.content[:12000] if ch.content else "",
                    },
                }
            )
        for i in range(0, len(vectors), 100):
            await index.upsert(vectors=vectors[i : i + 100])

        expires_at = datetime.now() + timedelta(hours=self.default_expiration_hours)
        existing_doc = await self.db_service.get_document(document_id)
        if existing_doc:
            await self.db_service.update_document(
                document_id,
                {
                    "processing_status": f"reingest_h{full_content_hash[:16]}",
                    "pinecone_index_name": index_name,
                    "pinecone_index_expires_at": expires_at.isoformat()
                    if isinstance(expires_at, datetime)
                    else expires_at,
                },
            )
        else:
            await self.db_service.create_document(
                user_id=user_id,
                name=file_name,
                file_name=file_name,
                file_type=file_type,
                file_size=len(file_content),
                content=text_content[:10000],
                pinecone_index_name=index_name,
                expires_at=expires_at,
                s3_key=s3_key,
                document_id=document_id,
                content_hash=full_content_hash,
                ingest_version=1,
            )

        try:
            db_rows: List[Dict[str, Any]] = []
            for ch in chunks:
                db_rows.append(
                    {
                        "id": str(uuid.uuid4()),
                        "document_id": document_id,
                        "chunk_index": ch.chunk_index,
                        "content": ch.content,
                        "metadata": ch.metadata,
                    }
                )
            await self.db_service.replace_document_chunks(document_id, db_rows)
        except Exception as db_err:
            logger.warning("document_chunks persist failed (Pinecone still updated): %s", db_err)

        await self.db_service.mark_document_processed(
            document_id=document_id,
            chunk_count=len(chunks),
            pinecone_index_name=index_name,
        )
        
        return {
            "document_id": document_id,
            "index_name": index_name,
            "chunk_count": len(chunks),
            "expires_at": expires_at.isoformat(),
            # "s3_key": s3_key,
            "status": "processed"
        }
    
    async def reingest_expired_document(self, document_id: str, user_id: str) -> Dict[str, Any]:
        """
        Re-ingest an expired document from S3
        
        Args:
            document_id: Document ID
            user_id: User ID
            
        Returns:
            Processing results
        """
        # Get document metadata
        doc = await self.db_service.get_document(document_id)
        if not doc:
            raise ValueError(f"Document {document_id} not found")
        
        # Check if file exists in S3
        if not self.s3_service.enabled:
            raise ValueError("S3 storage is not enabled. Cannot re-ingest expired document.")
        
        file_content = await self.s3_service.download_file(
            user_id=user_id,
            document_id=document_id,
            file_name=doc["file_name"]
        )
        
        if not file_content:
            raise ValueError(f"File not found in S3 for document {document_id}")
        
        # Process document again
        return await self.process_document(
            file_content=file_content,
            file_name=doc["file_name"],
            file_type=doc["file_type"],
            user_id=user_id,
            document_id=document_id
        )
    
    async def _extract_text(self, file_content: bytes, file_type: str, file_name: str) -> str:
        """Extract text from various file types"""
        if file_type == "application/pdf":
            return self._extract_pdf_text(file_content)
        elif file_type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            return self._extract_docx_text(file_content)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._extract_pptx_text(file_content)
        elif file_type.startswith("image/"):
            return await self._extract_image_text(file_content, file_type)
        elif file_type == "text/plain":
            return file_content.decode('utf-8', errors='ignore')
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            pdf_file = BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            raise ValueError(f"Error extracting PDF text: {str(e)}")
        return text
    
    def _extract_docx_text(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        text = ""
        try:
            docx_file = BytesIO(file_content)
            doc = DocxDocument(docx_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            raise ValueError(f"Error extracting DOCX text: {str(e)}")
        return text
    
    def _extract_pptx_text(self, file_content: bytes) -> str:
        """Extract text from PPTX file"""
        text = ""
        try:
            pptx_file = BytesIO(file_content)
            presentation = Presentation(pptx_file)
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    text += "\n".join(slide_text) + "\n\n"
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            raise ValueError(f"Error extracting PPTX text: {str(e)}")
        return text
    
    async def _extract_image_text(self, file_content: bytes, file_type: str) -> str:
        """Extract text from image using file parser agent"""
        try:
            parsed_result = await self.file_parser_agent.parse_binary_content(
                content_data=file_content,
                media_type=file_type,
                filename="image",
            )
            if parsed_result and hasattr(parsed_result, 'extracted_text'):
                return parsed_result.extracted_text
            elif isinstance(parsed_result, dict) and parsed_result.get('extracted_text'):
                return parsed_result['extracted_text']
            else:
                return "Image processed but no text extracted."
        except Exception as e:
            logger.error(f"Error extracting image text: {e}")
            return f"Error processing image: {str(e)}"
    
    def _is_text_nonsensical(self, text: str, threshold: float = 0.3) -> bool:
        """
        Check if extracted text is nonsensical.
        Uses heuristics: ratio of non-alphanumeric chars, word length distribution, etc.
        
        Args:
            text: Text to check
            threshold: Threshold for nonsensical score (0-1)
            
        Returns:
            True if text appears nonsensical
        """
        if not text or len(text.strip()) < 50:
            return True
        
        # Check ratio of non-alphanumeric characters
        alphanumeric_chars = sum(1 for c in text if c.isalnum() or c.isspace())
        total_chars = len(text)
        if total_chars == 0:
            return True
        
        alphanumeric_ratio = alphanumeric_chars / total_chars
        if alphanumeric_ratio < 0.5:
            return True
        
        # Check for excessive whitespace or special characters
        words = text.split()
        if len(words) < 10:
            return True
        
        # Check average word length (nonsensical text often has unusual patterns)
        avg_word_length = sum(len(word) for word in words) / len(words)
        if avg_word_length > 20 or avg_word_length < 2:
            return True
        
        # Check for repeated characters (common in OCR errors)
        repeated_patterns = len(re.findall(r'(.)\1{4,}', text))
        if repeated_patterns > len(text) * 0.1:
            return True
        
        return False
    
    def _chunk_document_streaming(
        self, 
        text: str, 
        chunk_size: int = 1000, 
        overlap: int = 200
    ) -> Generator[DocumentChunk, None, None]:
        """
        Split document into overlapping chunks using memory-efficient streaming.
        Yields chunks instead of storing all in memory.
        
        Args:
            text: Document text
            chunk_size: Words per chunk
            overlap: Overlapping words between chunks
            
        Yields:
            DocumentChunk objects
        """
        words = text.split()
        total_words = len(words)
        
        if total_words == 0:
            return
        
        chunk_index = 0
        i = 0
        
        while i < total_words:
            # Get chunk words
            chunk_words = words[i:min(i + chunk_size, total_words)]
            chunk_text = " ".join(chunk_words)
            
            # Yield chunk (memory efficient)
            yield DocumentChunk(
                chunk_index=chunk_index,
                content=chunk_text,
                metadata={
                    "start_word": i,
                    "end_word": min(i + chunk_size, total_words)
                }
            )
            
            # Move to next chunk with overlap
            i += chunk_size - overlap
            chunk_index += 1
            
            # Prevent infinite loop
            if i >= total_words:
                break
    
    async def _create_pinecone_index(self, index_name: str):
        """Create a Pinecone index for document storage"""
        try:
            # Check if index exists
            existing_indexes = await self.pinecone.list_indexes()
            if index_name in existing_indexes:
                return
            
            # Create index
            await self.pinecone.create_index(
                name=index_name,
                dimension=self.embedding_dimension,
                metric="cosine",
                spec=ServerlessSpec(
                    cloud="aws",
                    region=os.getenv("PINECONE_REGION", "us-east-1")
                ),
                deletion_protection="disabled"
            )
        except Exception as e:
            logger.error(f"Index creation error: {e}")
            raise
    
    async def search_documents(
        self,
        query: str,
        document_id: Optional[str] = None,
        top_k: int = 5
    ) -> List[Dict[str, Any]]:
        """
        Search documents using RAG.
        
        Args:
            query: Search query
            document_id: Optional document ID to search within
            top_k: Number of results to return
            
        Returns:
            List of relevant document chunks
        """
        if not document_id:
            raise ValueError("document_id is required for search")
        
        # Get index name from database
        doc = await self.db_service.get_document(document_id)
        if not doc:
            raise ValueError(f"Document {document_id} not found")
        
        index_name = doc.get("pinecone_index_name")
        if not index_name:
            # Check if document expired and can be re-ingested
            if self.s3_service.enabled and doc.get("s3_key"):
                logger.info(f"Document {document_id} expired, attempting re-ingestion from S3")
                try:
                    await self.reingest_expired_document(document_id, doc["user_id"])
                    doc = await self.db_service.get_document(document_id)
                    index_name = doc.get("pinecone_index_name")
                except Exception as e:
                    logger.error(f"Failed to re-ingest document: {e}")
                    raise ValueError(f"Document {document_id} has expired and cannot be re-ingested")
            else:
                raise ValueError(f"Document {document_id} has expired and S3 storage is not enabled")
        
        # Generate query embedding
        query_embedding = await self._generate_embedding(query)
        
        # Query Pinecone
        index = await self.pinecone.get_index(index_name)
        results = await index.query(
            vector=query_embedding,
            top_k=top_k,
            include_metadata=True,
            filter={"document_id": document_id}
        )
        rel_min = float(
            getattr(config, "DOCUMENT_GAME_RELEVANCE_THRESHOLD", 0.35) or 0.35
        )
        raw = [
            {
                "content": (match.metadata or {}).get("content", ""),
                "chunk_index": (match.metadata or {}).get("chunk_index", 0),
                "file_name": (match.metadata or {}).get("file_name", ""),
                "score": match.score,
            }
            for match in results.matches
        ]
        raw = filter_by_relevance(raw, threshold=rel_min)
        return dedup_text_snippets(raw, text_key="content")

    async def get_chunks_for_document(
        self, document_id: str, top_k: int = 100
    ) -> List[Dict[str, Any]]:
        """Get all chunks for a document (uses broad query for vector retrieval)."""
        return await self.search_documents(
            query="medical clinical health document content",
            document_id=document_id,
            top_k=top_k
        )

    async def get_adjacent_db_chunks(
        self,
        document_id: str,
        chunk_index: int,
        before: int = 1,
        after: int = 1,
    ) -> List[Dict[str, Any]]:
        """Pull neighboring persisted rows for wider context (small DB rows, not 5k-word blocks)."""
        return await self.db_service.get_document_chunks_adjacent(
            document_id, chunk_index, before=before, after=after
        )

# Global instance
_document_processor: Optional[DocumentProcessor] = None


def get_document_processor() -> DocumentProcessor:
    """Get or create global DocumentProcessor instance"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
